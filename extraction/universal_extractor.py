#!/usr/bin/env python3
"""
Universal document text extractor for government documents.

Handles: PDF, PowerPoint, Word, Excel, HTML, Images (OCR)

Usage:
    from extraction.universal_extractor import UniversalDocumentExtractor
    
    extractor = UniversalDocumentExtractor()
    result = extractor.extract_from_url("https://example.com/agenda.pdf")
    print(result['text'])
"""

import io
from pathlib import Path
from typing import Optional, Dict
import httpx
from loguru import logger

# PDF extraction
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None
    logger.warning("PDF support disabled. Install: pip install PyPDF2")

try:
    import pdfplumber
except ImportError:
    pdfplumber = None
    logger.debug("pdfplumber not available (optional)")

# PowerPoint extraction
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    logger.warning("PowerPoint support disabled. Install: pip install python-pptx")

# Word extraction
try:
    from docx import Document
except ImportError:
    Document = None
    logger.warning("Word support disabled. Install: pip install python-docx")

# Excel extraction
try:
    import pandas as pd
except ImportError:
    pd = None
    logger.warning("Excel support disabled. Install: pip install openpyxl pandas")

# HTML extraction
try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None
    logger.warning("HTML support disabled. Install: pip install beautifulsoup4")

# OCR extraction (for images/scanned PDFs)
try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None
    Image = None
    logger.debug("OCR support disabled (optional). Install: pip install pytesseract pillow")


class UniversalDocumentExtractor:
    """Extract text from any government document format."""
    
    def __init__(self):
        """Initialize extractor with HTTP client."""
        self.client = httpx.Client(timeout=30, follow_redirects=True)
    
    def extract_from_url(self, url: str) -> Dict[str, any]:
        """
        Download document from URL and extract text.
        
        Args:
            url: Document URL
            
        Returns:
            Dict with:
            - url: Source URL
            - format: File format (.pdf, .pptx, etc.)
            - text: Extracted text
            - file_size_kb: Size in KB
            - text_length: Length of extracted text
            - success: Whether extraction succeeded
        """
        logger.info(f"Downloading: {url}")
        
        try:
            # Download file
            response = self.client.get(url)
            response.raise_for_status()
            file_bytes = response.content
            
            # Detect format from URL or Content-Type
            file_ext = self._detect_format(url, response.headers.get('content-type', ''))
            
            logger.debug(f"Detected format: {file_ext}")
            
            # Extract based on format
            if file_ext == '.pdf':
                text = self.extract_pdf(file_bytes)
            elif file_ext in ['.ppt', '.pptx']:
                text = self.extract_powerpoint(file_bytes)
            elif file_ext in ['.doc', '.docx']:
                text = self.extract_word(file_bytes)
            elif file_ext in ['.xls', '.xlsx']:
                text = self.extract_excel(file_bytes)
            elif file_ext in ['.html', '.htm']:
                text = self.extract_html(file_bytes)
            elif file_ext in ['.jpg', '.jpeg', '.png', '.tiff', '.tif']:
                text = self.extract_image_ocr(file_bytes)
            else:
                logger.warning(f"Unknown format: {file_ext}")
                text = ""
            
            success = bool(text.strip())
            
            return {
                'url': url,
                'format': file_ext,
                'text': text,
                'file_size_kb': len(file_bytes) // 1024,
                'text_length': len(text),
                'success': success
            }
        
        except Exception as e:
            logger.error(f"Failed to extract from {url}: {e}")
            return {
                'url': url,
                'format': 'unknown',
                'text': '',
                'file_size_kb': 0,
                'text_length': 0,
                'success': False,
                'error': str(e)
            }
    
    def _detect_format(self, url: str, content_type: str) -> str:
        """Detect document format from URL or Content-Type."""
        
        # Try URL extension first
        url_lower = url.lower()
        for ext in ['.pdf', '.pptx', '.ppt', '.docx', '.doc', '.xlsx', '.xls', 
                    '.html', '.htm', '.jpg', '.jpeg', '.png', '.tiff', '.tif']:
            if ext in url_lower:
                return ext
        
        # Try Content-Type
        content_type_lower = content_type.lower()
        if 'pdf' in content_type_lower:
            return '.pdf'
        elif 'powerpoint' in content_type_lower or 'presentation' in content_type_lower:
            return '.pptx'
        elif 'word' in content_type_lower or 'msword' in content_type_lower:
            return '.docx'
        elif 'excel' in content_type_lower or 'spreadsheet' in content_type_lower:
            return '.xlsx'
        elif 'html' in content_type_lower:
            return '.html'
        elif 'image' in content_type_lower:
            return '.jpg'
        
        return '.unknown'
    
    def extract_pdf(self, file_bytes: bytes) -> str:
        """Extract text from PDF."""
        if PdfReader is None:
            logger.error("PyPDF2 not installed")
            return ""
        
        try:
            # Try PyPDF2 first (faster)
            pdf_reader = PdfReader(io.BytesIO(file_bytes))
            text = ""
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            
            # If no text extracted, might be scanned PDF
            if not text.strip() and pdfplumber:
                logger.info("PDF appears to be scanned, trying pdfplumber...")
                with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                    text = "\n".join(page.extract_text() or "" for page in pdf.pages)
            
            return text.strip()
        
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            return ""
    
    def extract_powerpoint(self, file_bytes: bytes) -> str:
        """Extract text from PowerPoint (.ppt, .pptx)."""
        if Presentation is None:
            logger.error("python-pptx not installed")
            return ""
        
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # Extract text from all shapes
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_parts.append(f"=== Slide {slide_num} ===")
                    text_parts.append("\n".join(slide_text))
                    text_parts.append("")
                
                # Extract speaker notes if available
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes:
                        text_parts.append(f"Notes: {notes}")
                        text_parts.append("")
            
            return "\n".join(text_parts).strip()
        
        except Exception as e:
            logger.error(f"PowerPoint extraction failed: {e}")
            return ""
    
    def extract_word(self, file_bytes: bytes) -> str:
        """Extract text from Word (.doc, .docx)."""
        if Document is None:
            logger.error("python-docx not installed")
            return ""
        
        try:
            doc = Document(io.BytesIO(file_bytes))
            
            text_parts = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            return "\n".join(text_parts).strip()
        
        except Exception as e:
            logger.error(f"Word extraction failed: {e}")
            return ""
    
    def extract_excel(self, file_bytes: bytes) -> str:
        """Extract text from Excel (.xls, .xlsx)."""
        if pd is None:
            logger.error("pandas/openpyxl not installed")
            return ""
        
        try:
            # Use pandas to read all sheets
            excel_file = io.BytesIO(file_bytes)
            all_sheets = pd.read_excel(excel_file, sheet_name=None, engine='openpyxl')
            
            text_parts = []
            for sheet_name, df in all_sheets.items():
                text_parts.append(f"=== Sheet: {sheet_name} ===")
                
                # Convert DataFrame to text
                text_parts.append(df.to_string(index=False))
                text_parts.append("")
            
            return "\n".join(text_parts).strip()
        
        except Exception as e:
            logger.error(f"Excel extraction failed: {e}")
            return ""
    
    def extract_html(self, file_bytes: bytes) -> str:
        """Extract text from HTML."""
        if BeautifulSoup is None:
            logger.error("BeautifulSoup not installed")
            return ""
        
        try:
            soup = BeautifulSoup(file_bytes, 'html.parser')
            
            # Remove script and style tags
            for script in soup(["script", "style", "nav", "header", "footer"]):
                script.decompose()
            
            # Get text
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return text.strip()
        
        except Exception as e:
            logger.error(f"HTML extraction failed: {e}")
            return ""
    
    def extract_image_ocr(self, file_bytes: bytes) -> str:
        """Extract text from image using OCR (for scanned documents)."""
        if pytesseract is None or Image is None:
            logger.error("pytesseract/PIL not installed")
            logger.info("Install: pip install pytesseract pillow")
            logger.info("Also install tesseract: sudo apt-get install tesseract-ocr")
            return ""
        
        try:
            image = Image.open(io.BytesIO(file_bytes))
            
            # Run OCR
            text = pytesseract.image_to_string(image)
            
            return text.strip()
        
        except Exception as e:
            logger.error(f"OCR extraction failed: {e}")
            logger.info("Make sure tesseract is installed: sudo apt-get install tesseract-ocr")
            return ""
    
    def close(self):
        """Close HTTP client."""
        self.client.close()
    
    def __enter__(self):
        """Context manager entry."""
        return self
    
    def __exit__(self, exc_type, exc_val, exc_tb):
        """Context manager exit."""
        self.close()


# Example usage and testing
if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python universal_extractor.py <url>")
        print("\nExample:")
        print("  python universal_extractor.py https://example.com/agenda.pdf")
        sys.exit(1)
    
    url = sys.argv[1]
    
    with UniversalDocumentExtractor() as extractor:
        result = extractor.extract_from_url(url)
        
        print(f"\n{'='*70}")
        print(f"URL: {result['url']}")
        print(f"Format: {result['format']}")
        print(f"File Size: {result['file_size_kb']} KB")
        print(f"Text Length: {result['text_length']} characters")
        print(f"Success: {result['success']}")
        print(f"{'='*70}\n")
        
        if result['success']:
            # Show first 500 characters
            preview = result['text'][:500]
            print("Preview:")
            print(preview)
            if len(result['text']) > 500:
                print("\n... (truncated)")
        else:
            print(f"Error: {result.get('error', 'Unknown error')}")
